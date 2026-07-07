import spire.presentation as sp
from datetime import datetime

class PPTEngine:
    def __init__(self):
        self.ppt = sp.Presentation()
        self.ppt.SlideSize.Type = sp.SlideSizeType.Screen16x9
        self.ppt.SlideSize.Orientation = sp.SlideOrienation.Landscape
        self.slide_width = 960
        self.slide_height = 540

    def add_slide(self):
        slide = self.ppt.Slides.Append()
        return slide

    def set_gradient_background(self, slide, color_start, color_end):
        bg = slide.SlideBackground
        bg.Type = sp.BackgroundType.Custom
        bg.Fill.FillType = sp.FillFormatType.Gradient
        grad = bg.Fill.Gradient
        grad.GradientStops.AppendByColor(0, sp.Color.FromArgb(255, *color_start))
        grad.GradientStops.AppendByColor(1, sp.Color.FromArgb(255, *color_end))

    def set_solid_background(self, slide, color):
        bg = slide.SlideBackground
        bg.Type = sp.BackgroundType.Custom
        bg.Fill.FillType = sp.FillFormatType.Solid
        bg.Fill.SolidColor.Color = sp.Color.FromArgb(255, *color)

    def create_card(self, slide, shape_type, x, y, width, height,
                    fill_color, alpha=200, shadow=True, dash_border=False):
        shape = slide.Shapes.AppendShape(
            shape_type,
            sp.RectangleF.FromLTRB(x, y, x + width, y + height)
        )
        r, g, b = fill_color
        if alpha <= 0:
            # 不可见卡片：完全透明且无边框
            shape.Fill.Visible = False
            shape.Line.Visible = False
        else:
            shape.Fill.FillType = sp.FillFormatType.Solid
            shape.Fill.SolidColor.Color = sp.Color.FromArgb(alpha, r, g, b)
            if dash_border:
                shape.Line.Visible = True
                shape.Line.Width = 1.5
                shape.Line.DashStyle = sp.LineDashStyleType.Dash
            else:
                shape.Line.Visible = False
        if shadow and alpha > 0 and hasattr(shape, 'Shadow'):
            try:
                shape.Shadow.ExternalShadowDistanceX = 6
                shape.Shadow.ExternalShadowDistanceY = 6
                shape.Shadow.ShadowColor = sp.Color.FromArgb(100, 0, 0, 0)
            except:
                pass
        return shape

    def add_text_to_shape(self, shape, text, font_name="微软雅黑",
                          font_size=18, bold=True, color=(255, 255, 255),
                          alignment=None):
        tf = shape.TextFrame
        tf.Text = text
        para = tf.Paragraphs[0]
        if alignment:
            para.Alignment = alignment
        else:
            para.Alignment = sp.TextAlignmentType.Center
        tr = para.TextRanges[0]
        tr.LatinFont = sp.TextFont(font_name)
        tr.EastAsianFont = sp.TextFont("微软雅黑")
        tr.FontHeight = font_size
        tr.IsBold = sp.TriState.TTrue if bold else sp.TriState.TFalse
        tr.Fill.FillType = sp.FillFormatType.Solid
        tr.Fill.SolidColor.Color = sp.Color.FromArgb(255, *color)

    def add_animation(self, slide, shape, anim_type):
        effect = slide.Timeline.MainSequence.AddEffect(shape, anim_type)
        return effect

    def set_transition(self, slide, transition_type):
        try:
            slide.SlideShowTransition.Type = transition_type
        except Exception as e:
            print(f"切换动画错误: {e}")

    def add_image(self, slide, image_path, x, y, width, height):
        # Spire.Presentation API 在不同版本中差异较大：
        # - 旧版本：slide.Shapes.AppendEmbedImage(sp.ShapeType.Image, path, rect)
        # - 新版本：slide.Shapes.AppendEmbedImageByPath(sp.ShapeType.Rectangle, path, rect)
        # 优先尝试新版本API，回退到旧版本
        rect = sp.RectangleF.FromLTRB(x, y, x + width, y + height)
        try:
            image = slide.Shapes.AppendEmbedImageByPath(sp.ShapeType.Rectangle, image_path, rect)
        except AttributeError:
            try:
                image = slide.Shapes.InsertEmbedImage(0, sp.ShapeType.Rectangle, rect, None)
                # 旧版本没有 Image shape type，InsertEmbedImage 需要 IImageData，跳过
                raise NotImplementedError
            except (AttributeError, NotImplementedError):
                # 最旧版本
                image = slide.Shapes.AppendEmbedImage(sp.ShapeType.Image, image_path, rect)
        try:
            image.Line.Visible = False
        except Exception:
            pass
        return image

    def add_decorative_shape(self, slide, shape_type, x, y, width, height, fill_color, alpha=80, rotation=0):
        shape = slide.Shapes.AppendShape(
            shape_type,
            sp.RectangleF.FromLTRB(x, y, x + width, y + height)
        )
        r, g, b = fill_color
        shape.Fill.FillType = sp.FillFormatType.Solid
        shape.Fill.SolidColor.Color = sp.Color.FromArgb(alpha, r, g, b)
        shape.Line.Visible = False
        if rotation != 0:
            shape.Rotation = rotation
        return shape

    def add_line(self, slide, x1, y1, x2, y2, color, width=2, dash_style=None):
        from spire.presentation import LineDashStyleType
        shape = slide.Shapes.AppendShape(
            sp.ShapeType.Line,
            sp.RectangleF.FromLTRB(x1, y1, x2, y2)
        )
        r, g, b = color
        shape.Line.FillType = sp.FillFormatType.Solid
        shape.Line.SolidFillColor.Color = sp.Color.FromArgb(255, r, g, b)
        shape.Line.Width = width
        if dash_style is not None:
            shape.Line.DashStyle = dash_style
        return shape

    def add_gradient_shape(self, slide, shape_type, x, y, width, height, color_start, color_end, alpha=200):
        shape = slide.Shapes.AppendShape(
            shape_type,
            sp.RectangleF.FromLTRB(x, y, x + width, y + height)
        )
        shape.Fill.FillType = sp.FillFormatType.Gradient
        grad = shape.Fill.Gradient
        grad.GradientStops.AppendByColor(0, sp.Color.FromArgb(alpha, *color_start))
        grad.GradientStops.AppendByColor(1, sp.Color.FromArgb(alpha, *color_end))
        shape.Line.Visible = False
        return shape

    def save(self, output_path, remove_borders=True):
        """保存PPT文件。

        Args:
            output_path: 输出路径
            remove_borders: 是否用python-pptx后处理移除所有shape边框。
                Spire.Presentation的Line.Visible=False无效（Line对象无Visible属性），
                边框只能通过python-pptx的shape.line.fill.background()移除。
                普通版（白底）必须设为True，专业版可设为False（装饰色块边框可能被误删）。
        """
        self.ppt.SaveToFile(output_path, sp.FileFormat.Pptx2016)
        import os
        size_kb = os.path.getsize(output_path) / 1024
        print(f"PPT已保存: {output_path}")
        print(f"文件大小: {size_kb:.1f} KB")
        print(f"总页数: {len(self.ppt.Slides)}")

        if remove_borders:
            self._remove_borders_postprocess(output_path)

        return output_path

    def _remove_borders_postprocess(self, pptx_path):
        """使用python-pptx后处理移除所有shape边框。

        Spire.Presentation的Line对象没有Visible属性，以下方法均无效：
        - shape.Line.Visible = False  （属性不存在）
        - shape.Line.Width = 0        （Spire忽略，仍渲染默认边框）
        - shape.Line.Style = TextLineStyle.none  （Style已是none但边框仍渲染）

        唯一有效方案：python-pptx的shape.line.fill.background()
        """
        try:
            from pptx import Presentation
            prs = Presentation(pptx_path)
            count = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not hasattr(shape, 'line'):
                        continue
                    try:
                        shape.line.fill.background()
                        count += 1
                    except Exception:
                        pass
            prs.save(pptx_path)
            print(f"边框后处理：已移除 {count} 个shape的边框")
        except ImportError:
            print("警告：python-pptx未安装，跳过边框移除后处理")
        except Exception as e:
            print(f"边框后处理失败: {e}")
